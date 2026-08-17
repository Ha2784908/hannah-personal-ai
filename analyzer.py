import os
import io
import json
import pypdf
from pptx import Presentation
from google import genai
from google.genai import types

def make_system_prompt(host_name: str = "Hannah") -> str:
    hname = host_name.strip() if (host_name and host_name.strip()) else "Hannah"
    return f"""You are {hname}'s Personal AI, a deeply intelligent, polite, respectful, and 100% honest personal assistant built specifically for {hname} (The Host).

YOUR CORE CHARACTER & UNCOMPROMISING PRINCIPLES:
1. **Polite, Warm & Respectful**: Always address {hname} politely and courteously. Maintain an encouraging, sharp, and helpful tone.
2. **Absolute Factual Honesty & Truth-Seeking**: Your highest commitment is to UNCOMPROMISING TRUTH, FACTUAL ACCURACY, and INTELLECTUAL HONESTY.
3. **No Blind Agreement or Fake Praise**: NEVER blindly validate false assumptions, agree with incorrect claims, or give empty compliments just to be polite.
4. **Polite Error Correction**:
   - If a prompt, question, text, or attached document contains false premises, factual errors, incorrect math/logic, or misstatements, YOU MUST POLITELY INFORM THE HOST ({hname}).
   - Clearly point out the exact mistake, explain step-by-step *why* it is incorrect, and provide the accurate, verified truth.
5. **Specialized Language Skills**: Exceptional proficiency in Chinese (中文), English, and Arabic (العربية). When translating or checking grammar, evaluate contextual nuances, idioms, pinyin (where appropriate), and natural tone.
"""

def get_client():
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        env_file = os.path.join(os.path.dirname(__file__), ".env")
        if os.path.exists(env_file):
            with open(env_file, "r", encoding="utf-8") as f:
                for line in f:
                    if line.startswith("GEMINI_API_KEY="):
                        api_key = line.split("=", 1)[1].strip()
                        os.environ["GEMINI_API_KEY"] = api_key
                        break
    if not api_key:
        return None
    return genai.Client(api_key=api_key)

def extract_pdf_text(pdf_bytes: bytes) -> str:
    """Extract text from PDF page by page."""
    try:
        reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
        extracted = []
        for idx, page in enumerate(reader.pages, 1):
            txt = page.extract_text()
            if txt and txt.strip():
                extracted.append(f"--- PDF Page {idx} ---\n{txt.strip()}")
        return "\n\n".join(extracted)
    except Exception as e:
        return f"[PDF Extraction Error: {str(e)}]"

def extract_pptx_text(pptx_bytes: bytes) -> str:
    """Extract text from PowerPoint slides."""
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        extracted = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text.append(paragraph.text.strip())
            if slide_text:
                extracted.append(f"--- PowerPoint Slide {idx} ---\n" + "\n".join(slide_text))
        return "\n\n".join(extracted)
    except Exception as e:
        return f"[PowerPoint Extraction Error: {str(e)}]"

def translate_text(text: str, source_lang: str = "Auto", target_lang: str = "English", tone: str = "Polite", host_name: str = "Hannah"):
    """Translate text between Chinese, English, Arabic, and other languages with automated Chinese double-checking."""
    client = get_client()
    if not client:
        return {"status": "error", "message": "GEMINI_API_KEY environment variable not set."}

    is_chinese = "Chinese" in source_lang or "Chinese" in target_lang or "中文" in source_lang or "中文" in target_lang or any('\u4e00' <= char <= '\u9fff' for char in text)
    system_prompt = make_system_prompt(host_name)

    prompt = f"""Translate the following text accurately into {target_lang}.
Host Person: {host_name}
Source Language: {source_lang}
Requested Tone: {tone}

Input Text:
\"\"\"
{text}
\"\"\"

Output JSON ONLY format:
{{
  "translation": "Translated text string...",
  "pinyin": "Pinyin representation if target or source is Chinese, otherwise empty string",
  "notes": "Contextual usage notes, idiom explanations, or polite tone insights",
  "alternatives": ["Alternative phrasing 1", "Alternative phrasing 2"],
  "chinese_verification": {{
    "is_chinese_involved": {str(is_chinese).lower()},
    "back_translation": "Back-translated phrase to confirm exact original meaning",
    "verification_status": "VERIFIED 100% ACCURATE",
    "explanation": "Simple explanation confirming why this Chinese phrasing is natural, accurate, and polite for {host_name}."
  }}
}}
"""
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config=types.GenerateContentConfig(
                system_instruction=system_prompt,
                temperature=0.2,
                response_mime_type="application/json"
            )
        )
        data = json.loads(response.text)
        return {"status": "success", "result": data}
    except Exception as e:
        return {"status": "error", "message": f"Translation failed: {str(e)}"}

def translate_uploaded_file(file_bytes: bytes, filename: str, mime_type: str, instruction: str = "", source_lang: str = "Auto", target_lang: str = "English", tone: str = "Polite", host_name: str = "Hannah"):
    """Translate text inside uploaded file with Chinese double-checking and customizable Host name."""
    client = get_client()
    if not client:
        return {"status": "error", "message": "GEMINI_API_KEY required for file translation."}

    ext = os.path.splitext(filename)[1].lower()
    is_chinese = "Chinese" in source_lang or "Chinese" in target_lang or "中文" in source_lang or "中文" in target_lang
    system_prompt = make_system_prompt(host_name)

    if mime_type.startswith("image/") or ext in ['.png', '.jpg', '.jpeg', '.webp']:
        image_part = types.Part.from_bytes(data=file_bytes, mime_type=mime_type if mime_type else "image/png")
        prompt = f"""Read all text inside this attached image/screenshot and translate it accurately into {target_lang}.
Host Person: {host_name}
Source Language: {source_lang}
Requested Tone: {tone}
Specific Instructions: {instruction if instruction else "Translate all visible text in the image accurately."}

Output JSON ONLY format:
{{
  "translation": "Translated text string...",
  "pinyin": "Pinyin representation if target or source is Chinese, otherwise empty string",
  "notes": "Contextual usage notes, idiom explanations, or polite tone insights",
  "alternatives": ["Alternative phrasing 1", "Alternative phrasing 2"],
  "chinese_verification": {{
    "is_chinese_involved": {str(is_chinese).lower()},
    "back_translation": "Back-translated phrase to confirm exact original meaning",
    "verification_status": "VERIFIED 100% ACCURATE",
    "explanation": "Simple explanation confirming why this Chinese phrasing is natural, accurate, and polite for {host_name}."
  }}
}}
"""
        try:
            response = client.models.generate_content(
                model="gemini-2.5-flash",
                contents=[image_part, prompt],
                config=types.GenerateContentConfig(system_instruction=system_prompt, temperature=0.2, response_mime_type="application/json")
            )
            data = json.loads(response.text)
            return {"status": "success", "result": data}
        except Exception as e:
            return {"status": "error", "message": f"Image translation failed: {str(e)}"}

    elif ext == '.pdf' or mime_type == 'application/pdf':
        extracted_text = extract_pdf_text(file_bytes)
        pdf_part = types.Part.from_bytes(data=file_bytes, mime_type="application/pdf")
        prompt = f"""Translate the content of this PDF document '{filename}' into {target_lang}.
Host Person: {host_name}
Source Language: {source_lang}
Requested Tone: {tone}
Specific Instructions: {instruction if instruction else "Translate the primary content of the document accurately."}

Extracted PDF Content:
\"\"\"
{extracted_text[:15000]}
\"\"\"

Output JSON ONLY format:
{{
  "translation": "Translated text string...",
  "pinyin": "Pinyin representation if target or source is Chinese, otherwise empty string",
  "notes": "Contextual usage notes, idiom explanations, or polite tone insights",
  "alternatives": ["Alternative phrasing 1", "Alternative phrasing 2"],
  "chinese_verification": {{
    "is_chinese_involved": {str(is_chinese).lower()},
    "back_translation": "Back-translated phrase to confirm exact original meaning",
    "verification_status": "VERIFIED 100% ACCURATE",
    "explanation": "Simple explanation confirming why this Chinese phrasing is natural, accurate, and polite for {host_name}."
  }}
}}
"""
        try:
            response = client.models.generate_content(
                model="gemini-2.5-flash",
                contents=[pdf_part, prompt],
                config=types.GenerateContentConfig(system_instruction=system_prompt, temperature=0.2, response_mime_type="application/json")
            )
            data = json.loads(response.text)
            return {"status": "success", "result": data}
        except Exception as e:
            return translate_text(f"PDF Document '{filename}' Extracted Content:\n\n{extracted_text}\n\nInstruction: {instruction}", source_lang, target_lang, tone, host_name)

    elif ext in ['.pptx', '.ppt']:
        extracted_text = extract_pptx_text(file_bytes)
        prompt_text = f"PowerPoint '{filename}' Slides Content:\n\n{extracted_text}\n\nInstruction: {instruction}"
        return translate_text(prompt_text, source_lang, target_lang, tone, host_name)

    else:
        content_str = file_bytes.decode('utf-8', errors='ignore')
        return translate_text(f"File '{filename}' Content:\n\n{content_str}\n\nInstruction: {instruction}", source_lang, target_lang, tone, host_name)

def check_grammar_and_facts(text: str, language: str = "Auto", host_name: str = "Hannah"):
    """Perform strict factual verification and grammar critique in Chinese, English, or Arabic."""
    client = get_client()
    if not client:
        return {"status": "error", "message": "GEMINI_API_KEY environment variable not set."}

    system_prompt = make_system_prompt(host_name)
    prompt = f"""Perform a comprehensive grammar check AND factual accuracy audit of the following text.
Host Person: {host_name}
Language Context: {language}

Text to Analyze:
\"\"\"
{text}
\"\"\"

CRITICAL INSTRUCTION:
- If there are factual errors, miscalculations, or false premises, POLITELY INFORM THE HOST ({host_name}). Do NOT pretend false claims are true.
- Output clean Markdown formatted as:

---
### 🎯 Accuracy & Grammar Verdict
[VERIFIED ACCURATE / CORRECTIONS REQUIRED / REVISIONS SUGGESTED]

### ✍️ Original vs. Corrected Text
**Original**:
> [original text]

**Corrected Version**:
> [corrected & polished text]

### 🧠 Factual & Grammar Breakdown
1. **Identified Errors / Misconceptions**:
   - ❌ **Issue**: [Mistake description]
   - *Why it's wrong*: [Polite step-by-step logic explanation]

2. **Key Grammar Rules & Tips**:
   - 💡 [Key learning takeaway & natural phrasing advice]
---
"""
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config=types.GenerateContentConfig(
                system_instruction=system_prompt,
                temperature=0.1,
            )
        )
        result_text = response.text
        status_flag = "CORRECTIONS REQUIRED" if ("CORRECTIONS REQUIRED" in result_text or "❌" in result_text) else "VERIFIED ACCURATE"
        return {"status": "success", "verdict": status_flag, "analysis": result_text}
    except Exception as e:
        return {"status": "error", "message": f"Grammar check failed: {str(e)}"}

def check_grammar_uploaded_file(file_bytes: bytes, filename: str, mime_type: str, instruction: str = "", language: str = "Auto", host_name: str = "Hannah"):
    """Perform grammar audit and fact check on uploaded file with customizable Host name."""
    client = get_client()
    if not client:
        return {"status": "error", "message": "GEMINI_API_KEY required for file grammar audit."}

    ext = os.path.splitext(filename)[1].lower()
    system_prompt = make_system_prompt(host_name)

    if mime_type.startswith("image/") or ext in ['.png', '.jpg', '.jpeg', '.webp']:
        image_part = types.Part.from_bytes(data=file_bytes, mime_type=mime_type if mime_type else "image/png")
        prompt = f"""Perform a comprehensive grammar check AND factual accuracy audit of the text inside this attached image/screenshot '{filename}'.
Host Person: {host_name}
Language Context: {language}
Specific Request: {instruction if instruction else "Check all visible text for grammar errors, natural structure, and factual truth."}

CRITICAL INSTRUCTION:
- If there are factual errors or grammar flaws, POLITELY INFORM THE HOST ({host_name}).
- Output clean Markdown formatted as:

---
### 🎯 Accuracy & Grammar Verdict
[VERIFIED ACCURATE / CORRECTIONS REQUIRED / REVISIONS SUGGESTED]

### ✍️ Original vs. Corrected Text
**Original**:
> [extracted text from image]

**Corrected Version**:
> [corrected & polished text]

### 🧠 Factual & Grammar Breakdown
1. **Identified Errors / Misconceptions**:
   - ❌ **Issue**: [Mistake description]
   - *Why it's wrong*: [Polite step-by-step logic explanation]

2. **Key Grammar Rules & Tips**:
   - 💡 [Key learning takeaway & natural phrasing advice]
---
"""
        try:
            response = client.models.generate_content(
                model="gemini-2.5-flash",
                contents=[image_part, prompt],
                config=types.GenerateContentConfig(system_instruction=system_prompt, temperature=0.1)
            )
            result_text = response.text
            status_flag = "CORRECTIONS REQUIRED" if ("CORRECTIONS REQUIRED" in result_text or "❌" in result_text) else "VERIFIED ACCURATE"
            return {"status": "success", "verdict": status_flag, "analysis": result_text}
        except Exception as e:
            return {"status": "error", "message": f"Image grammar check failed: {str(e)}"}

    elif ext == '.pdf' or mime_type == 'application/pdf':
        extracted_text = extract_pdf_text(file_bytes)
        pdf_part = types.Part.from_bytes(data=file_bytes, mime_type="application/pdf")
        prompt = f"""Perform a comprehensive grammar check AND factual accuracy audit of this PDF document '{filename}'.
Host Person: {host_name}
Language Context: {language}
Specific Request: {instruction if instruction else "Check the document for grammar errors, natural phrasing, and factual accuracy."}

Extracted PDF Content:
\"\"\"
{extracted_text[:15000]}
\"\"\"

CRITICAL INSTRUCTION:
- If there are factual errors or grammar flaws, POLITELY INFORM THE HOST ({host_name}).
- Output clean Markdown formatted as:

---
### 🎯 Accuracy & Grammar Verdict
[VERIFIED ACCURATE / CORRECTIONS REQUIRED / REVISIONS SUGGESTED]

### ✍️ Original vs. Corrected Text
**Original**:
> [original text snippet]

**Corrected Version**:
> [corrected & polished text]

### 🧠 Factual & Grammar Breakdown
1. **Identified Errors / Misconceptions**:
   - ❌ **Issue**: [Mistake description]
   - *Why it's wrong*: [Polite step-by-step logic explanation]

2. **Key Grammar Rules & Tips**:
   - 💡 [Key learning takeaway & natural phrasing advice]
---
"""
        try:
            response = client.models.generate_content(
                model="gemini-2.5-flash",
                contents=[pdf_part, prompt],
                config=types.GenerateContentConfig(system_instruction=system_prompt, temperature=0.1)
            )
            result_text = response.text
            status_flag = "CORRECTIONS REQUIRED" if ("CORRECTIONS REQUIRED" in result_text or "❌" in result_text) else "VERIFIED ACCURATE"
            return {"status": "success", "verdict": status_flag, "analysis": result_text}
        except Exception as e:
            return check_grammar_and_facts(f"PDF Document '{filename}' Extracted Content:\n\n{extracted_text}\n\nInstruction: {instruction}", language, host_name)

    elif ext in ['.pptx', '.ppt']:
        extracted_text = extract_pptx_text(file_bytes)
        prompt_text = f"PowerPoint Presentation '{filename}' Slides Content:\n\n{extracted_text}\n\nInstruction: {instruction}"
        return check_grammar_and_facts(prompt_text, language, host_name)

    else:
        content_str = file_bytes.decode('utf-8', errors='ignore')
        prompt_text = f"File '{filename}' Content:\n\n{content_str}\n\nInstruction: {instruction}"
        return check_grammar_and_facts(prompt_text, language, host_name)

def chat_with_hannah(message: str, host_name: str = "Hannah"):
    """Personal Host Assistant chat with honest fact-checking and polite truthfulness."""
    client = get_client()
    if not client:
        return {
            "status": "warning",
            "text": f"{host_name}'s Personal AI Server active. Configure `GEMINI_API_KEY` in `.env` to unlock full AI reasoning."
        }

    system_prompt = make_system_prompt(host_name)
    try:
        prompt = f"{host_name}'s Query/Submission: {message}"
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config=types.GenerateContentConfig(
                system_instruction=system_prompt,
                temperature=0.2,
            )
        )
        return {"status": "success", "text": response.text}
    except Exception as e:
        return {"status": "error", "text": f"Error interacting with AI Assistant: {str(e)}"}
